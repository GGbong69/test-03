"""기획서 PPTX 를 짓는다.

  python scripts/tools/make_deck.py

본문은 deck_pages.json 에서 오고, 부록 표는 data/*.csv 에서 직접 뽑는다.
용어를 바꾸려면 아래 TERMS 한 곳만 고친다 — 60장이 따라온다.
"""

import csv
import json
import re
import os
import sys

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
DATA = os.path.join(ROOT, "data")

try:                                   # cp949 콘솔에서도 한글이 깨지지 않게
    sys.stdout.reconfigure(encoding="utf-8")
except Exception:
    pass

# ── 용어 ──────────────────────────────────────────────────────────
# 왼쪽이 지금 본문에 적힌 말, 오른쪽이 화면에 나갈 말이다.
# 이름을 갈 때는 오른쪽만 고친다. 긴 말부터 바꾸므로 겹쳐도 안전하다.
TERMS = {
    "영역 강화 트랙": "영역 강화 트랙",   # 「트랙」 안의 「랙」을 막는 파수꾼
    "트랙": "트랙",
    "스티커 슬롯": "스티커 칸",
    "스티커 랙": "스티커 칸",
    "빈 랙": "빈 스티커 칸",
    "랙": "스티커 칸",
    "매대": "진열대",
    "정산 큐": "정산 순서",
    "탄창": "다트 묶음",
    "스티커": "스티커",
    "딱지": "딱지",
    "설비": "설비",
    "개조": "개조",
    "칩": "칩",
    "배수": "배수",
    "다트": "다트",
    "라운드": "라운드",
    "리그": "리그",
    "제약": "제약",
    "영역 강화": "영역 강화",
    "소비 아이템": "소비 아이템",
    "스타트 팩": "스타트 팩",
}
# 같은 말로 가는 짝도 남긴다 — 「트랙」처럼 긴 말이 먼저 걸려야
# 그 안의 짧은 말(「랙」)이 안 바뀐다.
_SUBS = dict(TERMS)
_PAT = re.compile("|".join(sorted((re.escape(k) for k in _SUBS), key=len, reverse=True)))     if _SUBS else None


def T(s):
    """용어를 한 번에 치환한다. 사슬로 물리지 않게 동시 치환이다."""
    s = "" if s is None else str(s)
    return _PAT.sub(lambda m: _SUBS[m.group(0)], s) if _PAT else s


def simul(obj, mapping):
    """자료 구조 전체의 문자열을 동시 치환한다."""
    pat = re.compile("|".join(sorted((re.escape(k) for k in mapping), key=len, reverse=True)))
    def go(o):
        if isinstance(o, str):
            return pat.sub(lambda m: mapping[m.group(0)], o)
        if isinstance(o, list):
            return [go(x) for x in o]
        if isinstance(o, dict):
            return {k: go(v) for k, v in o.items()}
        return o
    return go(obj)


# ── 서식 ──────────────────────────────────────────────────────────
F_BOLD = "Paperlogy 7 Bold"
F_SEMI = "Paperlogy 6 SemiBold"
F_BODY = "Paperlogy 4 Regular"

C_TEXT = RGBColor(0x17, 0x17, 0x17)
C_MUTE = RGBColor(0x62, 0x66, 0x6C)
C_ACC = RGBColor(0x16, 0x8D, 0xA8)
C_DARK = RGBColor(0x25, 0x28, 0x2D)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_RULE = RGBColor(0xD8, 0xD8, 0xD8)
C_BAND = RGBColor(0xF4, 0xF4, 0xF2)

L, R = 0.75, 12.58          # 좌우 여백
W = R - L                   # 본문 너비 11.83
Y_RULE = 1.25
Y_BODY = 1.46
Y_MAX = 6.92

W_BODY = 9.90        # 본문 폭. 오른쪽 2.7인치는 이미지 자리로 비운다
CW = W               # 지금 쪽의 본문 폭 (set_width 가 바꾼다)


def set_width(w):
    """이 쪽의 본문 폭을 정한다. 나란히 놓는 폭도 같이 따라간다."""
    global CW, HALF
    CW = w
    HALF = (w - 0.45) / 2


SZ_TITLE, SZ_LEAD, SZ_HEAD, SZ_BODY, SZ_SMALL = 20.25, 15.0, 13.5, 12.0, 10.5


def _apply_font(run, name, size, color):
    """라틴만이 아니라 한글(ea)·기타(cs) 자형까지 같은 글꼴로 박는다.
    font.name 만 쓰면 한글이 테마 기본 글꼴로 떨어진다."""
    run.font.name = name
    run.font.size = Pt(size)
    run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    latin = rPr.find(qn("a:latin"))
    idx = (list(rPr).index(latin) + 1) if latin is not None else len(rPr)
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            rPr.insert(idx, rPr.makeelement(qn(tag), {"typeface": name}))
            idx += 1
        else:
            el.set("typeface", name)


def _tf(shape):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def text(slide, x, y, w, h, s, size=SZ_BODY, font=F_BODY, color=C_TEXT,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=0.0):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = _tf(box)
    tf.vertical_anchor = anchor
    lines = T(s).split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = 1.25
        if space:
            p.space_before = Pt(space)
        run = p.add_run()
        run.text = line
        _apply_font(run, font, size, color)
    return box


def line(slide, x1, y1, x2, y2, color=C_RULE, width=0.75):
    ln = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    ln.line.color.rgb = color
    ln.line.width = Pt(width)
    return ln


def rect(slide, x, y, w, h, fill=None, edge=None, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if edge is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = edge
        sp.line.width = Pt(0.75)
    sp.shadow.inherit = False
    sp.text_frame.text = ""
    return sp


# ── 높이 어림 ──────────────────────────────────────────────────────
def _wrapped(s, chars):
    s = T(s)
    n = 1
    for part in s.split("\n"):
        n += max(0, (_wide_len(part) - 1) // chars)
    return n + len(s.split("\n")) - 1


def _wide_len(s):
    # 한글은 두 칸으로 센다
    return sum(2 if ord(c) > 0x2000 else 1 for c in s)


ROW_H = 0.265
HEAD_H = 0.36
# 폭 1인치가 담는 _wide_len 단위. 본문 11.5pt 한글 기준으로 잰 값이다.
UNITS_PER_IN = 12.6


def cap_of(width_in):
    return max(6, int((width_in - 0.30) * UNITS_PER_IN))


def table_geom(b, x=L, w=W):
    """표의 열 폭을 정한다. b['w'] 로 상대 비중을 직접 줄 수 있다."""
    n = max(1, len(b.get("cols") or []))
    weights = b.get("w")
    if not weights or len(weights) != n:
        weights = [1.0] * n
        if n >= 3:
            weights[0] = 0.78
            weights[-1] = 1.45
    tot = float(sum(weights))
    widths = [w * q / tot for q in weights]
    xs, acc = [], x
    for q in widths:
        xs.append(acc)
        acc += q
    return xs, widths


def row_height(row, widths, row_h=None):
    row_h = row_h or ROW_H
    lines = 1
    for ci, w in enumerate(widths):
        cell = row[ci] if ci < len(row) else ""
        lines = max(lines, _wrapped(cell, cap_of(w)))
    return row_h * lines


def label_w_for(w):
    return 2.20 if w > W * 0.6 else 1.45


def block_height(b, w=None):
    w = w or CW
    t = b.get("type")
    if t == "note":
        return 0.26 * _wrapped(b.get("text", ""), cap_of(w)) + 0.12
    if t == "flow":
        return 1.50
    if t == "chart":
        return (HEAD_H if b.get("heading") else 0.0) + b.get("h", 3.5) + 0.10
    head = HEAD_H if b.get("heading") else 0.0
    if t == "kv":
        rows = b.get("rows") or []
        cap = cap_of(w - label_w_for(w))
        return head + sum(ROW_H * _wrapped(r[1] if len(r) > 1 else "", cap) for r in rows) + 0.10
    if t == "split":
        half = (w - 0.45) / 2
        cap = cap_of(half - 1.45)
        a = sum(ROW_H * _wrapped(r[1] if len(r) > 1 else "", cap) for r in (b.get("rows") or []))
        c = sum(ROW_H * _wrapped(r[1] if len(r) > 1 else "", cap) for r in (b.get("rows2") or []))
        return HEAD_H + max(a, c) + 0.10
    if t == "steps":
        rows = b.get("rows") or []
        two = w > W * 0.62 and len(rows) > 4
        n = (len(rows) + 1) // 2 if two else len(rows)
        return head + n * 0.70 + 0.08
    if t == "table":
        _, widths = table_geom(b, L, w)
        rows = b.get("rows") or []
        rh = b.get("rh")
        return head + 0.30 + sum(row_height(r, widths, rh) for r in rows) + 0.10
    return 0.3


# ── 블록 그리기 ────────────────────────────────────────────────────
# 규칙 둘.
#   1. 표는 PPTX 의 진짜 표(a:tbl)로 넣는다. 도형 위에 텍스트 박스를 겹치지 않는다.
#   2. 한 박스에 들어갈 글은 한 박스에 쓴다. 줄마다 박스를 새로 만들지 않는다.
NO_STYLE = "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"      # 스타일 없음 · 눈금 없음


def _plain(table):
    """기본 표 스타일(파란 줄무늬·테두리)을 벗긴다."""
    pr = table._tbl.tblPr
    for k in ("firstRow", "bandRow", "lastRow", "firstCol", "lastCol", "bandCol"):
        pr.set(k, "0")
    el = pr.find(qn("a:tableStyleId"))
    if el is None:
        el = pr.makeelement(qn("a:tableStyleId"), {})
        pr.append(el)
    el.text = NO_STYLE


def _cell(c, s, size, font, color, bg=None, align=PP_ALIGN.LEFT):
    c.margin_left, c.margin_right = Inches(0.10), Inches(0.08)
    c.margin_top, c.margin_bottom = Inches(0.02), Inches(0.02)
    c.vertical_anchor = MSO_ANCHOR.MIDDLE
    if bg is None:
        c.fill.background()
    else:
        c.fill.solid()
        c.fill.fore_color.rgb = bg
    tf = c.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = T(s)
    _apply_font(r, font, size, color)


def make_table(slide, x, y, widths, heights, header, rows, band=True, size=None):
    """진짜 표 하나. header 가 None 이면 머리줄 없는 표다."""
    size = size or (SZ_BODY - 0.5)
    n = len(widths)
    nr = len(rows) + (1 if header else 0)
    gf = slide.shapes.add_table(nr, n, Inches(x), Inches(y),
                                Inches(sum(widths)), Inches(sum(heights)))
    t = gf.table
    _plain(t)
    for i, cw in enumerate(widths):
        t.columns[i].width = Inches(cw)
    for i, hh in enumerate(heights):
        t.rows[i].height = Inches(hh)
    off = 0
    if header:
        for ci in range(n):
            _cell(t.cell(0, ci), header[ci] if ci < len(header) else "",
                  SZ_SMALL, F_SEMI, C_WHITE, C_DARK)
        off = 1
    for ri, row in enumerate(rows):
        bg = C_BAND if (band and ri % 2 == 1) else None
        for ci in range(n):
            _cell(t.cell(off + ri, ci), row[ci] if ci < len(row) else "", size,
                  F_SEMI if ci == 0 else F_BODY, C_TEXT if ci == 0 else C_MUTE, bg)
    return y + sum(heights)


def draw_heading(slide, y, s, x=L, w=None):
    text(slide, x, y, w or CW, 0.30, s, size=SZ_HEAD, font=F_SEMI, color=C_ACC)
    return y + HEAD_H


def draw_kv(slide, y, b, x=L, w=None, label_w=None):
    """머리줄 없는 2열 표. 테두리도 줄무늬도 없다."""
    w = w or CW
    if label_w is None:
        label_w = label_w_for(w)
    if b.get("heading"):
        y = draw_heading(slide, y, b["heading"], x, w)
    rows = b.get("rows") or []
    if not rows:
        return y
    widths = [label_w, w - label_w]
    heights = [row_height(r, widths) for r in rows]
    return make_table(slide, x, y, widths, heights, None, rows,
                      band=False, size=SZ_BODY) + 0.10


def draw_split(slide, y, b, x=L, w=None):
    w = w or CW
    half = (w - 0.45) / 2
    ya = draw_kv(slide, y, {"heading": b.get("heading"), "rows": b.get("rows")},
                 x=x, w=half, label_w=1.45)
    yb = draw_kv(slide, y, {"heading": b.get("heading2"), "rows": b.get("rows2")},
                 x=x + half + 0.45, w=half, label_w=1.45)
    return max(ya, yb)


def draw_table(slide, y, b, x=L, w=None):
    w = w or CW
    cols = b.get("cols") or []
    rows = b.get("rows") or []
    if b.get("heading"):
        y = draw_heading(slide, y, b["heading"], x, w)
    if not rows:
        return y
    _, widths = table_geom(b, x, w)
    rh = b.get("rh")
    heights = [0.30] + [row_height(r, widths, rh) for r in rows]
    return make_table(slide, x, y, widths, heights, cols, rows,
                      size=b.get("size")) + 0.10


def draw_flow(slide, y, b, x=L, w=None):
    """둥근 상자 안에 글을 직접 넣는다 — 상자 따로 글 따로가 아니다."""
    w = w or CW
    items = [i for i in (b.get("items") or []) if i][:6]
    if not items:
        return y
    n = len(items)
    gap = 0.16
    bw = (w - gap * (n - 1)) / n
    bh = 1.00
    for i, s in enumerate(items):
        xx = x + i * (bw + gap)
        sp = rect(slide, xx, y, bw, bh, fill=None, edge=RGBColor(0xC9, 0xCC, 0xD0),
                  shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        tf = _tf(sp)
        tf.margin_left = tf.margin_right = Inches(0.06)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = T(s)
        _apply_font(r, F_SEMI, SZ_BODY, C_TEXT)
        if i < n - 1:
            line(slide, xx + bw + 0.02, y + bh / 2, xx + bw + gap - 0.02, y + bh / 2,
                 color=C_ACC, width=1.0)
    return y + bh + 0.22


def draw_steps(slide, y, b, x=L, w=None):
    """한 걸음이 텍스트 박스 하나. 번호 · 제목 · 설명을 한 박스에 담는다."""
    w = w or CW
    rows = b.get("rows") or []
    if b.get("heading"):
        y = draw_heading(slide, y, b["heading"], x, w)
    two = w > W * 0.62 and len(rows) > 4
    half = (len(rows) + 1) // 2 if two else len(rows)
    colw = (w - 0.45) / 2 if two else w
    for i, row in enumerate(rows):
        col, idx = ((0, i) if i < half else (1, i - half)) if two else (0, i)
        xx = x + col * (colw + 0.45)
        yy = y + idx * 0.70
        box = slide.shapes.add_textbox(Inches(xx), Inches(yy), Inches(colw), Inches(0.66))
        tf = _tf(box)
        p = tf.paragraphs[0]
        p.line_spacing = 1.2
        rn = p.add_run()
        rn.text = "%02d  " % (i + 1)
        _apply_font(rn, F_BOLD, SZ_BODY, C_ACC)
        rt = p.add_run()
        rt.text = T(row[0] if row else "")
        _apply_font(rt, F_SEMI, SZ_BODY, C_TEXT)
        p2 = tf.add_paragraph()
        p2.line_spacing = 1.2
        rd = p2.add_run()
        rd.text = T(row[1] if len(row) > 1 else "")
        _apply_font(rd, F_BODY, SZ_BODY - 0.5, C_MUTE)
        line(slide, xx, yy + 0.62, xx + colw, yy + 0.62)
    return y + half * 0.70 + 0.08


def draw_note(slide, y, b, x=L, w=None):
    w = w or CW
    h = 0.26 * _wrapped(b.get("text", ""), cap_of(w))
    text(slide, x, y, w, h, b.get("text", ""), size=SZ_BODY - 0.5, color=C_MUTE)
    return y + h + 0.12


def draw_chart(slide, y, b, x=L, w=None):
    """꺾은선 그래프. 목표 곡선처럼 수치보다 모양이 중요한 곳에 쓴다."""
    w = w or CW
    if b.get("heading"):
        y = draw_heading(slide, y, b["heading"], x, w)
    cd = CategoryChartData()
    cd.categories = [str(c) for c in b.get("cats", [])]
    for nm, vals in (b.get("series") or {}).items():
        cd.add_series(T(nm), vals)
    h = b.get("h", 3.5)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS,
                                Inches(x), Inches(y), Inches(w), Inches(h), cd)
    ch = gf.chart
    ch.has_title = False
    ch.has_legend = True
    ch.legend.position = XL_LEGEND_POSITION.BOTTOM
    ch.legend.include_in_layout = False
    ch.legend.font.size = Pt(10)
    ch.legend.font.name = F_BODY
    for ax in (ch.category_axis, ch.value_axis):
        ax.tick_labels.font.size = Pt(10)
        ax.tick_labels.font.name = F_BODY
    ch.value_axis.has_major_gridlines = True
    if b.get("log"):                       # 25 에서 32,000 까지 — 로그 축이라야 앞이 보인다
        vax = ch.value_axis._element
        lb = vax.makeelement(qn("c:logBase"), {"val": "10"})
        vax.insert(list(vax).index(vax.find(qn("c:scaling"))) if vax.find(qn("c:scaling")) is not None else 0, lb)
        sc = vax.find(qn("c:scaling"))
        if sc is not None:
            sc.insert(0, sc.makeelement(qn("c:logBase"), {"val": "10"}))
            vax.remove(lb)
    return y + h + 0.10


DRAW = {"chart": draw_chart, "kv": draw_kv, "split": draw_split, "table": draw_table,
        "flow": draw_flow, "steps": draw_steps, "note": draw_note}


# ── 슬라이드 ──────────────────────────────────────────────────────
def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def chrome(slide, title, page):
    text(slide, L, 0.44, W, 0.67, title, size=SZ_TITLE, font=F_BOLD, color=C_TEXT)
    line(slide, L, Y_RULE, R, Y_RULE, color=C_TEXT, width=1.0)
    text(slide, 11.9, 7.02, 0.68, 0.22, "%02d" % page, size=SZ_SMALL,
         color=C_MUTE, align=PP_ALIGN.RIGHT)


OVERFLOW = []
APPENDIX = []      # 부록 쪽 (번호, 제목) — 목차가 쓴다
TRIM = True          # 한 쪽을 넘기면 조금 덜어낸다. 그래도 넘치면 쪽을 나눈다
MAX_CUT = 40         # 한 쪽에 담기게 덜어낸다. 쪽을 쪼개지 않는다


def _bottom(related):
    return Y_MAX if related else Y_MAX + 0.34


HALF = (W - 0.45) / 2
GAP = 0.16


def fit_rows(b, avail, w=W):
    """표 b 의 행 중 avail 인치 안에 들어가는 개수."""
    _, widths = table_geom(b, L, w)
    room = avail - (HEAD_H if b.get("heading") else 0.0) - 0.30 - 0.10
    n = 0
    for r in b.get("rows") or []:
        h = row_height(r, widths)
        if room - h < 0:
            break
        room -= h
        n += 1
    return n


def _pairable(b):
    if b.get("type") in ("flow", "steps", "split"):
        return False
    if b.get("type") == "table" and len(b.get("cols") or []) > 4:
        return False
    return True


PAIR = True


def pack(blocks):
    """블록을 줄로 묶는다. 나란히 놓아 자리를 버는 둘은 한 줄에 넣는다."""
    lines, i = [], 0
    while i < len(blocks):
        a = blocks[i]
        b = blocks[i + 1] if i + 1 < len(blocks) else None
        if PAIR and b is not None and _pairable(a) and _pairable(b):
            pair = max(block_height(a, HALF), block_height(b, HALF))
            seq = block_height(a) + block_height(b) + GAP
            if pair <= seq * 0.82:
                lines.append([a, b])
                i += 2
                continue
        lines.append([a])
        i += 1
    return lines


def line_height(ln):
    if len(ln) == 2:
        return max(block_height(ln[0], HALF), block_height(ln[1], HALF))
    return block_height(ln[0])


def draw_line(slide, y, ln):
    if len(ln) == 2:
        y1 = DRAW[ln[0]["type"]](slide, y, ln[0], L, HALF)
        y2 = DRAW[ln[1]["type"]](slide, y, ln[1], L + HALF + 0.55, HALF)
        return max(y1, y2)
    return DRAW[ln[0]["type"]](slide, y, ln[0], L, W)


CUTS = []


def page_height(lead, lines):
    h = (0.34 * _wrapped(lead, cap_of(W)) + 0.22) if lead else 0.0
    return h + sum(line_height(ln) + 0.09 for ln in lines)


def fit_page(chapter, title, lead, blocks):
    global PAIR
    PAIR = chapter != "데이터"
    set_width(W if chapter == "데이터" else W_BODY)
    """한 쪽에 들어갈 때까지 덜어낸다. 뺀 것은 CUTS 에 남긴다."""
    blocks = [{**b,
               **({"rows": [list(r) for r in b["rows"]]} if b.get("rows") else {}),
               **({"rows2": [list(r) for r in b["rows2"]]} if b.get("rows2") else {})}
              for b in blocks]
    budget = Y_MAX - Y_BODY
    if page_height(lead, pack(blocks)) <= budget:
        return blocks

    # 1) note 를 뺀다 — 보조 문장이라 가장 먼저 버린다
    for b in list(blocks):
        if page_height(lead, pack(blocks)) <= budget:
            break
        if b.get("type") == "note":
            blocks.remove(b)
            CUTS.append("%s - %s · note 「%s」" % (chapter, title, (b.get("text") or "")[:44]))

    # 2) 가장 행이 많은 블록에서 끝 행부터 뺀다. 조금만 — 넘게 넘치면 쪽을 나눈다
    cut = 0
    while page_height(lead, pack(blocks)) > budget and cut < MAX_CUT:
        cand = [b for b in blocks if b.get("rows") and len(b["rows"]) > 2]
        if not cand:
            if len(blocks) > 1:                      # 마지막 수단 — 뒤 블록을 통째로
                gone = blocks.pop()
                CUTS.append("%s - %s · 블록 「%s」 통째" % (chapter, title,
                            gone.get("heading") or gone.get("type")))
                continue
            break
        cut += 1
        b = max(cand, key=lambda x: len(x["rows"]))
        gone = b["rows"].pop()
        CUTS.append("%s - %s · %s 「%s」" % (chapter, title,
                                            b.get("heading") or b.get("type"),
                                            " / ".join(map(str, gone))[:56]))
    return blocks


def content_slide(prs, page, chapter, title, lead, blocks, related):
    """한 페이지 분량을 그린다. 넘치면 「(이어서)」 쪽으로 넘긴다.
    다음에 쓸 쪽 번호를 돌려준다."""
    global PAIR
    PAIR = chapter != "데이터"
    set_width(W if chapter == "데이터" else W_BODY)
    if chapter == "데이터":
        APPENDIX.append(("%02d" % page, title))
    usable = [b for b in blocks if DRAW.get(b.get("type"))]
    if TRIM:
        usable = fit_page(chapter, title, lead, usable)
    queue = pack(usable)
    first = True
    while True:
        slide = new_slide(prs)
        chrome(slide, "%s - %s" % (chapter, title if first else title + " (이어서)"), page)
        y = Y_BODY
        if first and lead:
            h = 0.34 * _wrapped(lead, cap_of(W))
            text(slide, L, y, W, h, lead, size=SZ_LEAD, font=F_SEMI, color=C_TEXT)
            y += h + 0.22

        bottom = Y_MAX if related else Y_MAX + 0.34
        while queue:
            ln = queue[0]
            avail = bottom - y
            if line_height(ln) <= avail:
                y = draw_line(slide, y, ln) + 0.09
                queue.pop(0)
                continue
            if len(ln) == 1 and ln[0]["type"] == "table":
                b = ln[0]
                n = fit_rows(b, avail)
                if n >= 3 and len(b["rows"]) - n >= 2:
                    head = dict(b)
                    head["rows"] = b["rows"][:n]
                    draw_table(slide, y, head)
                    rest = dict(b)
                    rest["rows"] = b["rows"][n:]
                    queue[0] = [rest]
                    break
            if y > Y_BODY + 0.6:          # 이미 뭔가 놓였으면 다음 쪽으로 미룬다
                break
            y = draw_line(slide, y, ln) + 0.09    # 빈 쪽인데도 안 들어간다
            queue.pop(0)
            OVERFLOW.append("%02d %s - %s" % (page, chapter, title))

        if related:
            text(slide, L, Y_MAX + 0.08, W, 0.30,
                 "관련  " + "  ·  ".join(related), size=SZ_SMALL, font=F_SEMI, color=C_ACC)
        page += 1
        first = False
        if not queue:
            return page


def cover(prs):
    slide = new_slide(prs)
    text(slide, L, 1.72, W, 1.20, "HIGHTONE", size=64.5, font=F_BOLD, color=C_TEXT)
    text(slide, L, 3.02, 4.0, 0.44, "하이톤", size=23.25, font=F_SEMI, color=C_MUTE)
    line(slide, L, 3.82, 3.04, 3.82, color=C_TEXT, width=1.0)
    text(slide, L, 4.20, 5.0, 0.46, "게임 기획서", size=27.0, font=F_BOLD, color=C_TEXT)
    text(slide, L, 5.06, 6.0, 0.90,
         "다트 스코어어택 로그라이트 · PC\n기획 김건의", size=13.5, color=C_MUTE)
    text(slide, 11.9, 7.02, 0.68, 0.22, "01", size=SZ_SMALL, color=C_MUTE, align=PP_ALIGN.RIGHT)


def toc(prs, entries):
    """목차. 장에 번호를 매기고 쪽마다 1.1 을 단다 — 계층이 번호로 읽힌다.
    장은 단을 넘어 쪼개지 않는다. 한 단이 표 하나다."""
    slide = new_slide(prs)
    chrome(slide, "목차", 2)
    if not entries:
        return slide

    H_CH, H_IT, BUD = 0.36, 0.205, Y_MAX - Y_BODY
    blocks, cur = [], []
    for e in entries:                       # 장 단위로 묶는다
        if e[0] == "chapter" and cur:
            blocks.append(cur)
            cur = []
        cur.append(e)
    if cur:
        blocks.append(cur)

    cols, col, used = [], [], 0.0
    for b in blocks:                        # 넘치면 다음 단으로. 쪼개지 않는다
        h = H_CH + H_IT * (len(b) - 1)
        if col and used + h > BUD:
            cols.append(col)
            col, used = [], 0.0
        col += b
        used += h
    if col:
        cols.append(col)

    colw = (W - 2 * 0.5) / 3
    for ci, col in enumerate(cols[:3]):
        heights = [(H_CH if k == "chapter" else H_IT) for k, _, _ in col]
        x = L + ci * (colw + 0.5)
        gf = slide.shapes.add_table(len(col), 3, Inches(x), Inches(Y_BODY),
                                    Inches(colw), Inches(sum(heights)))
        t = gf.table
        _plain(t)
        t.columns[0].width = Inches(0.62)
        t.columns[1].width = Inches(colw - 1.50)
        t.columns[2].width = Inches(0.88)
        for k, h in enumerate(heights):
            t.rows[k].height = Inches(h)
        for k, (kind, no, name) in enumerate(col):
            if kind == "chapter":
                _cell(t.cell(k, 0), no.split("|")[0], SZ_HEAD - 0.5, F_BOLD, C_ACC)
                _cell(t.cell(k, 1), name, SZ_HEAD - 0.5, F_BOLD, C_TEXT)
                _cell(t.cell(k, 2), no.split("|")[1], SZ_SMALL, F_SEMI, C_MUTE,
                      align=PP_ALIGN.RIGHT)
            else:
                _cell(t.cell(k, 0), no.split("|")[0], SZ_BODY - 0.5, F_BODY, C_MUTE)
                _cell(t.cell(k, 1), name, SZ_BODY, F_BODY, C_TEXT)
                _cell(t.cell(k, 2), no.split("|")[1], SZ_BODY - 0.5, F_BODY, C_MUTE,
                      align=PP_ALIGN.RIGHT)
        y = Y_BODY
        for k, h in enumerate(heights):
            if col[k][0] == "chapter":
                line(slide, x, y + h - 0.05, x + colw, y + h - 0.05,
                     color=RGBColor(0xC0, 0xC4, 0xC9), width=0.9)
            y += h
    return slide


KB_PATH = os.path.join(ROOT, "docs", "export", "hightone_kb.json")
# 저장소는 아직 「앤티 > 판 > 라운드」로 적는다. 기획서는 「런 > 라운드 > 판」이다.
KB_RENAME = {"앤티": "라운드", "라운드": "판"}
_KB = None


def kb():
    global _KB
    if _KB is None:
        with open(KB_PATH, encoding="utf-8") as fh:
            _KB = simul(json.load(fh), KB_RENAME)
    return _KB


def tbl(name):
    return kb()["tables"].get(name, [])


def num(v, dash="—"):
    if v is None or v == "":
        return dash
    if isinstance(v, float) and v.is_integer():
        return str(int(v))
    return str(v)


def appendix_stickers(prs, page):
    global TRIM
    TRIM = False
    live = [r for r in tbl("items") if r.get("enabled")]
    rank = {"일반": 0, "희귀": 1, "레어": 2, "전설": 3}
    live.sort(key=lambda r: (rank.get(r.get("rarity_ko"), 9), r.get("cost") or 0, r["name"]))
    rows = []
    for r in live:
        eff = r.get("effect_ko") or ""
        gold = r.get("gold_ko") or ""
        if gold:
            eff = (eff + " · " + gold) if eff and eff != "—" else gold
        rows.append([r["name"], r.get("rarity_ko") or "—",
                     r.get("cond_ko") or "모든 다트", eff or "—",
                     "%sG" % num(r.get("cost"), "0")])
    per = 25
    npg = -(-len(rows) // per)
    size = -(-len(rows) // npg)
    for i in range(npg):
        chunk = rows[i * size:(i + 1) * size]
        content_slide(prs, page, "데이터",
                      "스티커 %d/%d" % (i + 1, npg),
                      "",
                      [{"type": "table",
                        "cols": ["이름", "희귀도", "조건", "효과", "가격"],
                        "w": [1.15, 0.62, 2.05, 2.0, 0.45],
                        "rh": 0.205, "size": 9.5, "rows": chunk}], [])
        page += 1
    return page, len(rows)


def appendix_tables(prs, page):
    global TRIM
    TRIM = True

    def emit(title, lead, blocks):
        nonlocal page
        page = content_slide(prs, page, "데이터", title, "", blocks, [])

    darts, mods = tbl("darts"), tbl("mods")
    emit("다트 · 개조", "다트 %d종 · 개조 %d종" % (len(darts), len(mods)),
         [{"type": "table", "rh": 0.22, "size": 10.0, "heading": "다트 — 표준 1발을 대체한다",
           "cols": ["이름", "게이지", "배수", "효과", "가격"],
           "w": [1.0, 0.7, 0.5, 3.2, 0.6],
           "rows": [[d["name"], "×%s" % num(d["gauge"]), num(d["mult"], "0"),
                     d.get("desc_ko") or "—",
                     "%sG" % num(d["cost"]) if d["cost"] else "기본"] for d in darts]},
          {"type": "table", "rh": 0.22, "size": 10.0, "heading": "개조 — 판의 기하를 바꾼다 · 런 동안 유지",
           "cols": ["이름", "축", "효과", "배타", "가격"],
           "w": [1.0, 0.8, 2.8, 0.8, 0.6],
           "rows": [[m["name"], m["axis"], m.get("desc_ko") or "—",
                     m.get("excl") or "—", "%sG" % num(m["cost"])] for m in mods]}])

    md = tbl("modifiers")
    emit("보스 제약", "%d종 · 보스 판 입장 전 3장 중 1장을 고른다" % len(md),
         [{"type": "table", "rh": 0.22, "size": 10.0,
           "cols": ["이름", "무엇을 건드리나", "효과"],
           "w": [1.0, 1.1, 3.4],
           "rows": [[m["name"], m["axis"], m.get("desc_ko") or "—"] for m in md]}])

    vo, tg = tbl("vouchers"), tbl("tags")
    emit("설비 · 딱지", "설비 %d종 · 딱지 %d종" % (len(vo), len(tg)),
         [{"type": "table", "rh": 0.22, "size": 10.0, "heading": "설비 — 상점 왼쪽 벽, 라운드마다 하나",
           "cols": ["이름", "효과", "가격", "라운드", "선행"],
           "w": [1.1, 2.6, 0.5, 0.5, 1.0],
           "rows": [[v["name"], v.get("desc_ko") or "—", "%sG" % num(v["cost"]),
                     num(v["min_ante"]), v.get("prereq") or "—"] for v in vo]},
          {"type": "table", "rh": 0.22, "size": 10.0, "heading": "딱지 — 소판·대판을 건너뛸 때만 받는다",
           "cols": ["이름", "효과", "시점", "라운드"],
           "w": [1.1, 2.9, 0.9, 0.5],
           "rows": [[t["name"], t.get("desc_ko") or "—", t["when"], num(t["min_ante"])]
                    for t in tg]}])

    st, pk = tbl("stakes"), tbl("packs")
    emit("리그 · 스타트 팩", "리그 %d단 · 팩 %d종" % (len(st), len(pk)),
         [{"type": "table", "rh": 0.22, "size": 10.0,
           "heading": "리그 — 앞 단을 완주하면 다음 단이 열린다",
           "cols": ["리그", "곡선", "작은 판 골드", "봉인", "다트", "진열대 가격", "삭음", "유지비"],
           "rows": [[s_["name"], s_["curve"], num(s_["reward_small"], "0"),
                     num(s_["seal_items"], "0"), num(s_["darts_add"], "0"),
                     "×%s" % num(s_["shop_cost_mul"]), num(s_["perish"], "0"),
                     num(s_["rent"], "0")] for s_ in st]},
          {"type": "table", "rh": 0.22, "size": 10.0, "heading": "스타트 팩",
           "cols": ["팩", "다트", "시작 골드", "스티커 칸", "소비 칸", "대가와 변화"],
           "w": [1.0, 0.55, 0.7, 0.7, 0.6, 3.4],
           "rows": [[q["name"], num(q["darts_add"], "0"), num(q["gold_add"], "0"),
                     num(q["item_slots"]), num(q["cons_slots"]),
                     (q.get("desc_ko") or q.get("desc") or "—")] for q in pk]}])

    au = tbl("area_upgrades")
    tracks = {}
    for r in au:
        tracks.setdefault(r["name"], []).append(r)
    tn = kb()["tuning"]["values"]
    keys = ["start_gold", "clear_gold", "gold_per_dart", "interest_per", "interest_max",
            "sell_div", "sell_min", "free_rerolls", "reroll_base", "reroll_step",
            "max_items", "darts_base", "cons_slots", "curve_first", "curve_last"]
    emit("영역 강화 · 튜닝 상수",
         "강화 %d트랙 · 레벨당 증분" % len(tracks),
         [{"type": "table", "rh": 0.22, "size": 10.0, "heading": "영역 강화 — 런 동안 유지",
           "cols": ["트랙", "레벨", "레벨당 점수", "레벨당 배수"],
           "rows": [[k, str(len(v)),
                     " · ".join(num(x["add_score"], "0") for x in v),
                     " · ".join(num(x["add_mult"], "—") for x in v)]
                    for k, v in tracks.items()]},
          {"type": "table", "rh": 0.22, "size": 10.0, "heading": "튜닝 상수",
           "cols": ["키", "값", "키", "값"],
           "rows": [[keys[i], num(tn.get(keys[i])),
                     keys[i + 1] if i + 1 < len(keys) else "",
                     num(tn.get(keys[i + 1])) if i + 1 < len(keys) else ""]
                    for i in range(0, len(keys), 2)]}])

    return page


# ── 조립 ──────────────────────────────────────────────────────────
SECTIONS = [
    ("01", "용어", "", "점수와 판 · 모으는 것 · 런 바깥"),
    ("02", "기획 의도", "", "한 줄 정의 · 왜 다트판인가 · 무엇을 덜어냈나 · 실력 × 빌드 · 타깃"),
    ("03", "시스템", "", "런 · 판 · 다트 · 조준 · 점수 · 스티커 · 조건 · 개조 · 강화 · 제약"),
    ("04", "재미 요소", "", "코어 루프 · 조준 · 조합 · 발동 순서 · 선택의 저울"),
    ("05", "보상 요소", "", "골드 · 이자 · 상점 · 딜러 · 설비 · 딱지 · 팩 · 리그"),
    ("06", "UI/UX", "", "레이아웃 · 조작 · 화면 전이 · 상태별 처리"),
    ("07", "데이터", "", "스티커 전종 · 다트 · 개조 · 제약 · 리그 · 튜닝 상수"),
]


ORDER = ["용어", "기획 의도", "시스템", "재미 요소", "보상 요소", "UI/UX"]


def _run(pages, toc_rows=None):
    """덱을 한 번 짓는다. spans 가 없으면 목차의 쪽 범위를 비워 둔다."""
    prs = Presentation()
    prs.slide_width = Inches(13.3333)
    prs.slide_height = Inches(7.5)

    cover(prs)
    toc(prs, toc_rows or [])

    APPENDIX.clear()
    rows, page, cn = [], 3, 0
    for name in ORDER:
        got = [q for q in pages if q.get("chapter") == name]
        if not got:
            continue
        cn += 1
        rows.append(["chapter", "%d|" % cn, name])
        head, first, it = len(rows) - 1, page, 0
        for p in got:
            it += 1
            rows.append(["page", "%d.%d|%02d" % (cn, it, page), p["title"]])
            page = content_slide(prs, page, p["chapter"], p["title"], p.get("lead", ""),
                                 p.get("blocks", []), p.get("related", []))
        rows[head][1] += "%d–%d" % (first, page - 1)

    start = page
    page, n_stick = appendix_stickers(prs, page)
    page = appendix_tables(prs, page)
    cn += 1
    rows.append(["chapter", "%d|%d–%d" % (cn, start, page - 1), "데이터"])
    seen = set()
    it = 0
    for no, nm in APPENDIX:
        base = nm.split(" ")[0] if nm.startswith("스티커") else nm
        if base in seen:
            continue
        seen.add(base)
        it += 1
        rows.append(["page", "%d.%d|%s" % (cn, it, no),
                     "스티커 전종 99장" if base == "스티커" else nm])
    return prs, rows, page - 1, n_stick


def build(pages_json, out_path):
    with open(pages_json, encoding="utf-8") as fh:
        pages = json.load(fh)["pages"]

    _, rows, _, _ = _run(pages)             # 1차 — 쪽 번호를 잰다
    OVERFLOW.clear()
    CUTS.clear()
    prs, _, total, n_stick = _run(pages, rows)    # 2차 — 목차를 채워 짓는다
    prs.save(out_path)
    return total, n_stick


if __name__ == "__main__":
    src = sys.argv[1] if len(sys.argv) > 1 else os.path.join(ROOT, "docs", "deck_pages.json")
    dst = sys.argv[2] if len(sys.argv) > 2 else os.path.join(ROOT, "docs", "HIGHTONE_기획서.pptx")
    total, ns = build(src, dst)
    print("%d장 · 스티커 %d장 · %s" % (total, ns, dst))
    if CUTS:
        rep = os.path.join(ROOT, "docs", "deck_cuts.txt")
        with open(rep, "w", encoding="utf-8") as fh:
            fh.write("\n".join(CUTS) + "\n")
        print("지면에 안 들어가 뺀 것 %d줄 · docs/deck_cuts.txt" % len(CUTS))
    if OVERFLOW:
        print("아직 넘치는 쪽 %d개: %s" % (len(OVERFLOW), " / ".join(OVERFLOW[:6])))
